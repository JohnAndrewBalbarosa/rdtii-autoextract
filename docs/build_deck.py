"""Generate a formal Zetarix presentation (PPTX) with speaker notes.

Run: python docs/build_deck.py
Output: docs/RDTII_AutoExtract.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ----------------------------------------------------------------------------
# Design tokens — disciplined, formal palette (deep navy / slate / restrained accent)
# ----------------------------------------------------------------------------
NAVY = RGBColor(0x0E, 0x2A, 0x47)        # primary dark
SLATE = RGBColor(0x33, 0x44, 0x55)       # body text
MUTED = RGBColor(0x6B, 0x7B, 0x8C)       # secondary text
ACCENT = RGBColor(0x1F, 0x6F, 0xB2)      # restrained blue accent
ACCENT_2 = RGBColor(0x2E, 0x8B, 0x6F)    # teal-green for secondary marks
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)       # surface
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD6, 0xDD, 0xE5)        # hairline rule

FONT_HEAD = "Georgia"          # serif for headings — editorial, formal
FONT_BODY = "Calibri"          # clean sans for body

SW, SH = Inches(13.333), Inches(7.5)     # 16:9

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, fill, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = line_w or Pt(1)
    shp.shadow.inherit = False
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def setpara(p, text, size, color, bold=False, font=FONT_BODY, align=PP_ALIGN.LEFT,
            space_after=6, space_before=0, italic=False):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return p


def bullet(tf, text, size=16, color=SLATE, bold=False, level=0, space_after=8,
           first=False, marker="—", marker_color=ACCENT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    p.space_after = Pt(space_after)
    p.alignment = PP_ALIGN.LEFT
    # marker run
    rm = p.add_run()
    rm.text = marker + "  "
    rm.font.size = Pt(size)
    rm.font.bold = True
    rm.font.color.rgb = marker_color
    rm.font.name = FONT_BODY
    # text run
    rt = p.add_run()
    rt.text = text
    rt.font.size = Pt(size)
    rt.font.bold = bold
    rt.font.color.rgb = color
    rt.font.name = FONT_BODY
    return p


def notes(slide, script):
    slide.notes_slide.notes_text_frame.text = script


def header(slide, kicker, title, n):
    """Standard formal content-slide header with kicker, title, footer."""
    # top accent bar
    rect(slide, 0, 0, SW, Inches(0.14), NAVY)
    # kicker
    tb, tf = textbox(slide, Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.35))
    setpara(tf.paragraphs[0], kicker.upper(), 12, ACCENT, bold=True, font=FONT_BODY,
            space_after=0)
    # title
    tb, tf = textbox(slide, Inches(0.7), Inches(0.78), Inches(11.9), Inches(0.9))
    setpara(tf.paragraphs[0], title, 30, NAVY, bold=True, font=FONT_HEAD, space_after=0)
    # hairline
    rect(slide, Inches(0.7), Inches(1.62), Inches(11.93), Pt(1.4), RULE)
    footer(slide, n)


def footer(slide, n):
    tb, tf = textbox(slide, Inches(0.7), Inches(7.05), Inches(8), Inches(0.3))
    setpara(tf.paragraphs[0], "Zetarix  ·  UN ESCAP & KMITL Global Hackathon 2026",
            9, MUTED, font=FONT_BODY, space_after=0)
    tb, tf = textbox(slide, Inches(12.0), Inches(7.05), Inches(0.9), Inches(0.3))
    setpara(tf.paragraphs[0], str(n), 9, MUTED, font=FONT_BODY, align=PP_ALIGN.RIGHT,
            space_after=0)


def chip(slide, x, y, text, w=Inches(2.0), fill=LIGHT, fg=NAVY):
    h = Inches(0.42)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = RULE
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    setpara(tf.paragraphs[0], text, 11, fg, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, space_after=0)
    return box


# ============================================================================
# SLIDE 1 — Title
# ============================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
# left accent column
rect(s, 0, 0, Inches(0.35), SH, ACCENT)
# kicker
tb, tf = textbox(s, Inches(0.9), Inches(1.5), Inches(11), Inches(0.4))
setpara(tf.paragraphs[0],
        "GLOBAL HACKATHON ON AI FOR DIGITAL TRADE REGULATORY ANALYSIS  ·  UN ESCAP & KMITL  ·  2026",
        13, RGBColor(0x9F, 0xC4, 0xE6), bold=True, font=FONT_BODY, space_after=0)
# title
tb, tf = textbox(s, Inches(0.9), Inches(2.3), Inches(11.6), Inches(2.0))
setpara(tf.paragraphs[0], "Zetarix", 60, WHITE, bold=True, font=FONT_HEAD,
        space_after=6)
p = tf.add_paragraph()
setpara(p, "An Open-Source, Model-Agnostic AI Pipeline for Automated Digital Trade "
           "Regulatory Analysis Across Asia-Pacific Jurisdictions",
        20, RGBColor(0xD7, 0xE4, 0xF1), font=FONT_BODY, space_after=6)
p = tf.add_paragraph()
setpara(p, "RDTII  —  Regional Digital Trade Integration Index (UN ESCAP)",
        15, RGBColor(0x9F, 0xC4, 0xE6), bold=True, font=FONT_BODY, space_after=0)
# rule
rect(s, Inches(0.9), Inches(4.75), Inches(6), Pt(2), ACCENT)
# subline
tb, tf = textbox(s, Inches(0.9), Inches(5.0), Inches(11), Inches(1.2))
setpara(tf.paragraphs[0],
        "Automating ~80% of the RDTII discover → describe workflow — "
        "Pillar 6 (Cross-Border Data Flows) and Pillar 7 (Domestic Data Protection) — "
        "at article-level granularity, with transparent human review for the final ~20%.",
        15, RGBColor(0xBF, 0xD2, 0xE5), font=FONT_BODY, space_after=0)
tb, tf = textbox(s, Inches(0.9), Inches(6.5), Inches(11), Inches(0.5))
setpara(tf.paragraphs[0], "Licensed under Apache 2.0", 13, RGBColor(0x9F, 0xC4, 0xE6),
        bold=True, font=FONT_BODY, space_after=0)

notes(s,
"Good day, and thank you for the opportunity to present. This presentation introduces "
"Zetarix, an open-source, model-agnostic artificial intelligence pipeline that we "
"have developed for the Global Hackathon on AI for Digital Trade Regulatory Analysis, "
"jointly organised by UN ESCAP and KMITL.\n\n"
"The central objective of the project is to automate approximately eighty percent of the "
"RDTII regulatory workflow — specifically the discovery and description stages — while "
"deliberately preserving a transparent human review step for the remaining twenty percent. "
"Our scope addresses the two mandatory pillars: Pillar 6, Cross-Border Data Flows, and "
"Pillar 7, Domestic Data Protection, at article-level granularity. The entire system is "
"released under the Apache 2.0 licence to ensure it is freely reusable by the policy "
"community. Over the next slides I will outline the problem, our objectives, the system "
"architecture, and the measurable outcomes we are targeting.")


# ============================================================================
# SLIDE 2 — Problem Statement
# ============================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
header(s, "The Problem", "A Manual Workflow That Does Not Scale", 2)

tb, tf = textbox(s, Inches(0.7), Inches(1.95), Inches(6.3), Inches(4.6))
bullet(tf, "ESCAP researchers must manually search government portals, retrieve "
           "regulatory documents, read legal text, and extract structured data.",
       size=16, first=True, space_after=14)
bullet(tf, "This is performed at article-level granularity, across many Asia-Pacific "
           "jurisdictions and multiple languages.", size=16, space_after=14)
bullet(tf, "The process is time-intensive, difficult to scale, and dependent on "
           "researcher availability and language coverage.", size=16, space_after=14)
bullet(tf, "Documents outside the existing ESCAP RDTII database are easily missed.",
       size=16, space_after=14, marker_color=ACCENT_2)

# right stat panel
panel = rect(s, Inches(7.5), Inches(2.0), Inches(5.1), Inches(4.3), LIGHT)
tb, tf = textbox(s, Inches(7.85), Inches(2.35), Inches(4.4), Inches(3.7))
setpara(tf.paragraphs[0], "THE BURDEN TODAY", 12, MUTED, bold=True, space_after=14)
for big, small in [("100%", "manual effort — search, read, extract, map"),
                   ("Many", "jurisdictions and languages to cover"),
                   ("Per-article", "granularity required for every document")]:
    p = tf.add_paragraph()
    setpara(p, big, 34, ACCENT, bold=True, font=FONT_HEAD, space_after=0)
    p2 = tf.add_paragraph()
    setpara(p2, small, 13, SLATE, space_after=16)

notes(s,
"Allow me to begin with the problem we set out to solve. At present, the RDTII workflow is "
"almost entirely manual. ESCAP researchers must search government portals by hand, retrieve "
"regulatory documents, read through dense legal text, and then extract structured data — and "
"they must do this at article-level granularity, not merely at the document level.\n\n"
"This work spans numerous Asia-Pacific jurisdictions and several languages, which makes it "
"time-intensive and very difficult to scale. Throughput is bounded by researcher availability "
"and by the languages each researcher can read. Furthermore, any regulatory document that is "
"not already in the ESCAP database is easily overlooked. The figures on the right summarise "
"the burden: the effort is essentially one hundred percent manual, it must cover many "
"jurisdictions and languages, and it demands per-article precision. This is precisely the "
"bottleneck our pipeline is designed to relieve.")


# ============================================================================
# SLIDE 3 — Objectives
# ============================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
header(s, "Our Goal", "Project Objectives", 3)

objectives = [
    ("Automate discovery", "Locate regulatory documents beyond the existing ESCAP RDTII database."),
    ("Article-level extraction", "Capture all six mandatory fields: title, last update, URL, scope, provisions, impact."),
    ("Traceable mapping", "Map provisions to Pillar 6 and 7 sub-indicators with verifiable citations."),
    ("Self-hostable & open", "Deliver an open-source tool deployable without proprietary infrastructure."),
    ("Model-agnostic", "Make every AI component swappable to open-weight models (e.g. Llama 3) via config."),
    ("Human-first review", "Provide a non-technical UI where any mapping can be verified or rejected in seconds."),
]
x0, y0 = Inches(0.7), Inches(2.0)
cw, ch = Inches(3.86), Inches(2.1)
gx, gy = Inches(0.18), Inches(0.2)
for i, (t, d) in enumerate(objectives):
    col = i % 3
    row = i // 3
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    card = rect(s, x, y, cw, ch, LIGHT)
    bar = rect(s, x, y, Inches(0.09), ch, ACCENT if row == 0 else ACCENT_2)
    tb, tf = textbox(s, x + Inches(0.3), y + Inches(0.25), cw - Inches(0.55), ch - Inches(0.45))
    setpara(tf.paragraphs[0], f"0{i+1}", 14, MUTED, bold=True, font=FONT_HEAD, space_after=4)
    p = tf.add_paragraph(); setpara(p, t, 18, NAVY, bold=True, font=FONT_HEAD, space_after=6)
    p = tf.add_paragraph(); setpara(p, d, 13, SLATE, space_after=0)

notes(s,
"Our objectives follow directly from that problem. First, we aim to automate document "
"discovery so that the system can locate relevant regulations even when they lie outside the "
"existing ESCAP database. Second, we target genuine article-level extraction, capturing all "
"six mandatory fields — title, last update, URL, scope, provisions, and impact — for every "
"article.\n\n"
"Third, every extracted provision must be mapped to the correct Pillar 6 or Pillar 7 "
"sub-indicator, and each mapping must carry a verifiable citation back to the source. Fourth, "
"the tool must be self-hostable and fully open-source, with no dependency on proprietary "
"infrastructure in production. Fifth — and this is central to our design philosophy — every "
"AI component must be model-agnostic and swappable to open-weight models through configuration "
"alone. Finally, we commit to a human-first review experience: a non-technical interface in "
"which any mapping can be verified, edited, or rejected within seconds.")


# ============================================================================
# SLIDE 4 — Architecture / Model-agnostic
# ============================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
header(s, "Design Philosophy", "Ports & Adapters: A Model-Agnostic Core", 4)

tb, tf = textbox(s, Inches(0.7), Inches(1.95), Inches(5.9), Inches(4.5))
bullet(tf, "Built on ports & adapters (hexagonal) architecture.", size=16, bold=True,
       first=True, space_after=12)
bullet(tf, "The core domain depends only on interfaces — never on a concrete model.",
       size=15, space_after=12)
bullet(tf, "Every AI model (LLM, OCR, captioner, embeddings, graph) is a swappable "
           "suggestion, changed via config — no domain edits.", size=15, space_after=12)
bullet(tf, "Development can use a paid API; production resolves to open-weight "
           "(Llama 3.1 via Ollama) with a single config change.", size=15, space_after=12,
       marker_color=ACCENT_2)
bullet(tf, "Swappability is explicitly scored: 20 points in Stage 1 plus 20 points in "
           "Stage 3.", size=15, bold=True, space_after=0, marker_color=ACCENT_2)

# right diagram: core hexagon-ish with adapters
dx = Inches(7.0)
core = rect(s, dx + Inches(1.55), Inches(3.4), Inches(2.3), Inches(1.2), NAVY)
tb, tf = textbox(s, dx + Inches(1.55), Inches(3.55), Inches(2.3), Inches(0.9), anchor=MSO_ANCHOR.MIDDLE)
setpara(tf.paragraphs[0], "CORE DOMAIN", 13, WHITE, bold=True, font=FONT_BODY,
        align=PP_ALIGN.CENTER, space_after=2)
p = tf.add_paragraph(); setpara(p, "ports only", 11, RGBColor(0x9F,0xC4,0xE6),
        align=PP_ALIGN.CENTER, italic=True, space_after=0)
adapters = ["LLM", "OCR", "Captioner", "Embeddings", "Vector store", "Graph"]
positions = [(0.2, 2.1), (2.9, 2.1), (5.0, 3.5), (0.0, 4.9), (2.9, 4.9), (5.0, 5.7)]
for name, (ax, ay) in zip(adapters, positions):
    chip(s, dx + Inches(ax), Inches(ay), name, w=Inches(1.55), fill=LIGHT, fg=NAVY)

notes(s,
"This slide explains the single most important architectural decision in the project: the "
"system is model-agnostic by design. We build on a ports and adapters — also called "
"hexagonal — architecture. The core domain, which contains the business logic of the RDTII "
"workflow, depends only on interfaces, which we call ports. It never references a concrete "
"model directly.\n\n"
"Every artificial intelligence component — the large language model, the OCR engine, the "
"captioner, the embedding model, and the graph library — is implemented as an adapter behind "
"one of those ports. As a result, each model is a swappable suggestion rather than a hardcoded "
"commitment. We can run a paid commercial API during development for convenience, and then, "
"for production, resolve every call to an open-weight model such as Llama 3.1 served through "
"Ollama — changing only configuration, with no edits to the domain code. This is not merely an "
"engineering preference: the hackathon rubric explicitly rewards swappability, awarding twenty "
"points in Stage 1 and a further twenty points in Stage 3. The diagram on the right shows the "
"core surrounded by interchangeable adapters.")


# ============================================================================
# SLIDE 5 — Pipeline (5 stages)
# ============================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
header(s, "How It Works", "The Five-Stage Pipeline", 5)

stages = [
    ("1", "Discovery & Retrieval", "Crawl national portals, gazettes, and ministry sites; capture full provenance per document.", ACCENT),
    ("2", "OCR & Captioning", "Convert any format to clean text at <5% CER; SigLIP-style captions become the basis for tags.", ACCENT),
    ("3", "Tagging & Extraction", "Multi-label tag each article; RAG extracts all 6 fields with citations and confidence scores.", ACCENT),
    ("4", "Human Review", "Non-technical UI: Accept, Reject, or Edit each mapping in seconds, linked to the source.", ACCENT_2),
    ("5", "Concept Graph", "Connect tagged provisions into a weighted, pruned, navigable graph — the core contribution.", ACCENT_2),
]
x0, y = Inches(0.7), Inches(2.3)
cw = Inches(2.28)
gx = Inches(0.13)
for i, (n, t, d, col) in enumerate(stages):
    x = x0 + i * (cw + gx)
    card = rect(s, x, y, cw, Inches(3.4), LIGHT)
    rect(s, x, y, cw, Inches(0.7), col)
    tb, tf = textbox(s, x, y + Inches(0.08), cw, Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
    setpara(tf.paragraphs[0], "STAGE " + n, 13, WHITE, bold=True, align=PP_ALIGN.CENTER, space_after=0)
    tb, tf = textbox(s, x + Inches(0.2), y + Inches(0.95), cw - Inches(0.4), Inches(2.3))
    setpara(tf.paragraphs[0], t, 15, NAVY, bold=True, font=FONT_HEAD, space_after=8)
    p = tf.add_paragraph(); setpara(p, d, 12, SLATE, space_after=0)
    if i < 4:
        ar = textbox(s, x + cw - Inches(0.02), y + Inches(1.4), Inches(0.2), Inches(0.4),
                     anchor=MSO_ANCHOR.MIDDLE)[1]
        setpara(ar.paragraphs[0], "›", 24, MUTED, bold=True, align=PP_ALIGN.CENTER, space_after=0)

tb, tf = textbox(s, Inches(0.7), Inches(6.05), Inches(11.9), Inches(0.5))
setpara(tf.paragraphs[0],
        "discover  →  describe  (Stages 1–4)        +        connect  (Stage 5)",
        14, MUTED, italic=True, font=FONT_BODY, align=PP_ALIGN.CENTER, space_after=0)

notes(s,
"Here is the pipeline as a whole. It runs in five stages. The first four stages cover the "
"discover-to-describe arc that the hackathon requires, and the fifth stage adds our "
"distinctive contribution.\n\n"
"In Stage 1, Discovery and Retrieval, a crawler targets national legislation portals, official "
"gazettes, and ministry websites, recording full provenance for every document. In Stage 2, "
"OCR and Captioning, we convert any document format into clean, machine-readable text at below "
"five percent character error rate; importantly, a SigLIP-style caption is produced for each "
"section and becomes the basis for tagging, so the tags are not held hostage to OCR noise. "
"In Stage 3, Tagging and Extraction, each article is multi-label tagged, and a "
"retrieval-augmented generation step extracts all six mandatory fields together with "
"citations and a confidence score. Stage 4 is the Human Review interface, where a "
"non-technical researcher can accept, reject, or edit any mapping in seconds, with each record "
"linked back to its source. Finally, Stage 5 connects the tagged provisions into a weighted, "
"pruned, navigable concept graph — which I will describe in more detail shortly.")


# ============================================================================
# SLIDE 6 — Concept Graph (core contribution)
# ============================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
header(s, "Core Contribution", "Stage 5 — The Concept Graph", 6)

tb, tf = textbox(s, Inches(0.7), Inches(1.95), Inches(6.4), Inches(4.6))
bullet(tf, "Edges weight tag overlap (IDF-weighted) and embedding cosine similarity "
           "between provisions.", size=15, first=True, space_after=12)
bullet(tf, "Pseudo-deterministic threshold pruning: edges below a calibrated θ are "
           "removed; results reproduce given fixed seeds and θ.", size=15, space_after=12)
bullet(tf, "Community detection (NetworkX / Louvain — BSD) produces Obsidian-style "
           "topical groupings.", size=15, space_after=12)
bullet(tf, "Formal Concept Analysis builds a generality hierarchy; weighted PageRank "
           "ranks entry points and orders within levels.", size=15, space_after=12,
       marker_color=ACCENT_2)
bullet(tf, "Result: broad entry-point categories → ranked specific sub-topics — a "
           "navigable map of related provisions.", size=15, bold=True, space_after=0,
       marker_color=ACCENT_2)

# right: simple node-graph illustration (upper) + formal model panel (lower)
gx, gy = Inches(7.5), Inches(2.1)
gw, gh = Inches(5.1), Inches(2.95)
rect(s, gx, gy, gw, gh, LIGHT)
import math
cx, cy = gx + gw/2, gy + gh/2
# central node
def node(xc, yc, r, color, label, fs=11, fg=WHITE):
    nd = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(xc - r)), Emu(int(yc - r)),
                            Emu(int(2*r)), Emu(int(2*r)))
    nd.fill.solid(); nd.fill.fore_color.rgb = color
    nd.line.color.rgb = WHITE; nd.line.width = Pt(1.5)
    nd.shadow.inherit = False
    tf2 = nd.text_frame; tf2.word_wrap = True
    setpara(tf2.paragraphs[0], label, fs, fg, bold=True, align=PP_ALIGN.CENTER, space_after=0)
    return (xc, yc)

def edge(p1, p2):
    ln = s.shapes.add_connector(2, Emu(int(p1[0])), Emu(int(p1[1])), Emu(int(p2[0])), Emu(int(p2[1])))
    ln.line.color.rgb = RULE; ln.line.width = Pt(1.5)
    ln.shadow.inherit = False
    # send to back
    sp = ln._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(3, sp)

R = Inches(0.62)
satellites = []
for k in range(6):
    ang = math.radians(60 * k - 90)
    sxc = cx + math.cos(ang) * Inches(1.55)
    syc = cy + math.sin(ang) * Inches(0.92)
    satellites.append((sxc, syc))
center = (cx, cy)
for sat in satellites:
    edge(center, sat)
# a couple cross edges
edge(satellites[0], satellites[1])
edge(satellites[3], satellites[4])
labels = ["Transfer", "Consent", "Localization", "Breach", "DPA", "Adequacy"]
cols = [ACCENT_2, ACCENT_2, ACCENT, ACCENT, ACCENT_2, ACCENT]
for (sxc, syc), lb, c in zip(satellites, labels, cols):
    node(sxc, syc, Inches(0.44), c, lb, fs=9)
node(cx, cy, Inches(0.58), NAVY, "Data\nFlows", fs=11)

# --- Formal model panel: the equations behind the graph ---
ex, ey = Inches(7.5), Inches(5.25)
ew, eh = Inches(5.1), Inches(1.7)
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ex, ey, ew, eh)
panel.fill.solid(); panel.fill.fore_color.rgb = NAVY
panel.line.fill.background(); panel.shadow.inherit = False
etb, etf = textbox(s, ex + Inches(0.25), ey + Inches(0.12), ew - Inches(0.5), eh - Inches(0.24))
setpara(etf.paragraphs[0], "FORMAL MODEL", 11, RGBColor(0x9F, 0xC4, 0xE6),
        bold=True, font=FONT_BODY, space_after=8)
eqs = [
    ("Edge weight", "w(i, j) = α · tag(i, j) + β · cos(i, j)"),
    ("Pruning", "keep (i, j)  ⟺  w(i, j) ≥ θ"),
    ("Ranking", "PR = d · Mᵀ·PR + (1 − d) / N"),
]
for name, formula in eqs:
    p = etf.add_paragraph()
    r1 = p.add_run(); r1.text = name + ":  "
    r1.font.size = Pt(11); r1.font.bold = True
    r1.font.name = FONT_BODY; r1.font.color.rgb = RGBColor(0xD7, 0xE4, 0xF1)
    r2 = p.add_run(); r2.text = formula
    r2.font.size = Pt(12.5); r2.font.bold = False
    r2.font.name = "Cambria Math"; r2.font.color.rgb = WHITE
    p.space_after = Pt(3)
# legend for the symbols
lp = etf.add_paragraph()
setpara(lp, "tag = IDF-weighted tag overlap · cos = embedding cosine similarity · "
            "θ = pruning threshold · M = adjacency matrix",
        9, RGBColor(0x9F, 0xC4, 0xE6), italic=True, font=FONT_BODY, space_after=0)

notes(s,
"Stage 5 is what we consider the core architectural contribution of the project, because it "
"goes beyond extraction and turns isolated records into knowledge. After tagging, each article "
"becomes a node. We then draw weighted edges between nodes based on two signals: the overlap "
"of their tags, weighted by inverse document frequency, and the cosine similarity of their "
"embeddings.\n\n"
"To keep the graph meaningful and reproducible, we apply pseudo-deterministic threshold "
"pruning: every edge below a calibrated threshold, theta, is removed, and because we fix the "
"random seeds and theta, the same inputs always produce the same graph. We then run community "
"detection using NetworkX and the Louvain method — deliberately chosen because it is "
"BSD-licensed and therefore compatible with our Apache repository. On top of this we apply "
"Formal Concept Analysis to build a generality hierarchy, where provisions with more tags are "
"more specific, and weighted PageRank to rank entry points. The outcome is a navigable map: "
"the reader starts from broad entry-point categories and drills down into ranked, specific "
"sub-topics, with related provisions cross-linked automatically.")


# ============================================================================
# SLIDE 7 — Tech Stack
# ============================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
header(s, "Implementation", "Technology Stack — Every Layer Behind a Port", 7)

rows = [
    ("Layer (port)", "Suggested adapter (dev)", "Open-weight target", True),
    ("LLM  ·  LLMProvider", "Claude / GPT", "Llama 3.1 8B/70B via Ollama", False),
    ("Captioning  ·  Captioner", "Qwen2-VL / ColQwen2 (Apache)", "Same; OCR-text fallback", False),
    ("Embeddings  ·  Tagger", "BGE-M3 (multilingual, MIT)", "Same", False),
    ("OCR  ·  OCREngine", "Tesseract 5 + OpenCV / PaddleOCR", "Same", False),
    ("Vector store", "pgvector / Chroma (self-hostable)", "Same", False),
    ("Graph + ranking", "NetworkX (BSD) + FCA concepts (MIT)", "Same", False),
    ("Backend / Frontend", "FastAPI  ·  React / Next.js", "Same", False),
    ("Deployment", "Docker + docker-compose", "Same", False),
]
x0, y0 = Inches(0.7), Inches(1.95)
col_w = [Inches(3.7), Inches(4.7), Inches(3.53)]
rh = Inches(0.5)
for r, (a, b, c, hdr) in enumerate(rows):
    y = y0 + r * rh
    fill = NAVY if hdr else (LIGHT if r % 2 else WHITE)
    rect(s, x0, y, sum(col_w, Emu(0)), rh, fill,
         line=RULE, line_w=Pt(0.5))
    cells = [a, b, c]
    cx = x0
    for ci, (txt, w) in enumerate(zip(cells, col_w)):
        tb, tf = textbox(s, cx + Inches(0.18), y, w - Inches(0.3), rh, anchor=MSO_ANCHOR.MIDDLE)
        color = WHITE if hdr else (NAVY if ci == 0 else SLATE)
        setpara(tf.paragraphs[0], txt, 12.5 if not hdr else 12, color,
                bold=(hdr or ci == 0), font=FONT_BODY, space_after=0)
        cx = cx + w

tb, tf = textbox(s, Inches(0.7), Inches(6.6), Inches(11.9), Inches(0.5))
setpara(tf.paragraphs[0],
        "All components are Apache-2.0-compatible. Named models are reference suggestions, not commitments.",
        12, MUTED, italic=True, font=FONT_BODY, space_after=0)

notes(s,
"This table sets out the technology stack, but the way to read it is important. Each row is a "
"layer that sits behind a port — an interface. The middle column lists the adapter we suggest "
"for development, and the right column lists the open-weight target we resolve to in "
"production. The point is that these are reference suggestions, not hard commitments; any of "
"them can be replaced through configuration.\n\n"
"For the language model we may use Claude or GPT in development and Llama 3.1 in production. "
"For captioning we use Qwen2-VL or ColQwen2, both Apache-licensed, with an OCR-text fallback. "
"Embeddings use the multilingual BGE-M3 model. OCR is handled by Tesseract 5 with OpenCV, "
"falling back to PaddleOCR for Southeast Asian scripts. The vector store is pgvector or Chroma, "
"both self-hostable. The graph layer uses NetworkX and an MIT-licensed Formal Concept Analysis "
"library. Backend is FastAPI, frontend is Next.js, and the whole system ships with Docker "
"Compose. Critically, every component is Apache-2.0-compatible, which is a licensing "
"requirement of the competition, and each reused component will be disclosed in the Technical "
"Memo with its licence verified.")


# ============================================================================
# SLIDE 8 — Coverage, Cost & Sustainability
# ============================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
header(s, "Evidence & Viability", "Coverage, Cost, and Sustainability", 8)

cards = [
    ("Coverage", ACCENT, [
        "Round 1: evaluated on 5–10 documents from 3 provided countries.",
        "At least one non-English jurisdiction handled via the translation pipeline.",
        "Finale-ready for 10 assigned countries — no retraining required.",
    ]),
    ("Cost", ACCENT_2, [
        "Open-weight (local Llama 3.1 8B): ≈ USD 0.00 in API fees per 50-page document.",
        "Cloud GPU compute: ≈ USD 0.05–0.10 per document at batch rates.",
        "API configuration: ≈ USD 0.10–0.25 per document.",
    ]),
    ("Sustainability", NAVY, [
        "Fully self-hostable via Docker on commodity hardware.",
        "No proprietary API dependency in production.",
        "Modular: OCR, LLM, and vector DB upgrade independently.",
    ]),
]
x0, y0 = Inches(0.7), Inches(2.0)
cw = Inches(3.86)
gx = Inches(0.18)
for i, (title, col, items) in enumerate(cards):
    x = x0 + i * (cw + gx)
    rect(s, x, y0, cw, Inches(4.35), LIGHT)
    rect(s, x, y0, cw, Inches(0.62), col)
    tb, tf = textbox(s, x, y0 + Inches(0.06), cw, Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
    setpara(tf.paragraphs[0], title.upper(), 15, WHITE, bold=True, font=FONT_BODY,
            align=PP_ALIGN.CENTER, space_after=0)
    tb, tf = textbox(s, x + Inches(0.28), y0 + Inches(0.85), cw - Inches(0.5), Inches(3.3))
    for j, it in enumerate(items):
        bullet(tf, it, size=13, space_after=12, first=(j == 0), marker_color=col)

notes(s,
"Beyond the design, the project must be credible in practice, so this slide addresses "
"coverage, cost, and sustainability. On coverage, in Round 1 the tool will be evaluated "
"against five to ten documents from three provided countries, including at least one "
"non-English jurisdiction handled by our translation pipeline. The architecture is designed to "
"scale to the ten assigned countries at the finale without any retraining.\n\n"
"On cost, the figures are deliberately conservative. In the open-weight configuration, running "
"Llama 3.1 8B locally, the API fee for a fifty-page document is effectively zero; the only cost "
"is cloud GPU compute, approximately five to ten US cents per document at batch rates. If a "
"commercial API is used instead, the cost is roughly ten to twenty-five cents per document. On "
"sustainability, the entire pipeline is self-hostable through Docker on commodity hardware, "
"with no proprietary dependency in production, and its modular design means the OCR engine, the "
"language model, and the vector database can each be upgraded independently over time. Exact "
"benchmarks will be reported in the Technical Memo after empirical testing.")


# ============================================================================
# SLIDE 9 — Milestones / Roadmap
# ============================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, WHITE)
header(s, "The Plan", "Key Milestones", 9)

milestones = [
    ("25 May 2026", "Application submitted — CVs, Concept Video, Technical Memo, Declaration"),
    ("31 May 2026", "Shortlist announcement"),
    ("5–10 Jun 2026", "ESCAP & KMITL training workshops"),
    ("20 Jul 2026", "Round 1 submission — tool evaluated on 3 provided countries"),
    ("31 Jul 2026", "Hybrid pitch"),
    ("30 Sep 2026", "Final submission"),
    ("15 Oct 2026", "Award ceremony, Bangkok"),
]
# vertical timeline
lx = Inches(2.7)
y0 = Inches(2.05)
step = Inches(0.66)
rect(s, lx, y0, Pt(2.5), step * (len(milestones) - 1) + Inches(0.1), RULE)
for i, (date, ev) in enumerate(milestones):
    y = y0 + i * step
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, lx - Inches(0.08), y - Inches(0.02),
                             Inches(0.22), Inches(0.22))
    col = ACCENT_2 if i == 0 else ACCENT
    dot.fill.solid(); dot.fill.fore_color.rgb = col
    dot.line.fill.background(); dot.shadow.inherit = False
    tb, tf = textbox(s, Inches(0.7), y - Inches(0.06), Inches(1.9), Inches(0.4))
    setpara(tf.paragraphs[0], date, 13, NAVY, bold=True, font=FONT_BODY,
            align=PP_ALIGN.RIGHT, space_after=0)
    tb, tf = textbox(s, lx + Inches(0.35), y - Inches(0.08), Inches(9.0), Inches(0.5))
    setpara(tf.paragraphs[0], ev, 14, SLATE, font=FONT_BODY, space_after=0)

notes(s,
"This slide lays out the timeline through to the award ceremony. The immediate milestone is the "
"application submission on the twenty-fifth of May 2026, which comprises our CVs, the concept "
"video, the Technical Memo, and the signed declaration. The shortlist is announced at the end "
"of May, followed by ESCAP and KMITL training workshops in early June.\n\n"
"The first major technical checkpoint is the Round 1 submission on the twentieth of July, where "
"the tool is evaluated on three provided countries, followed by the hybrid pitch at the end of "
"that month. The final submission is due on the thirtieth of September, and the award ceremony "
"takes place in Bangkok on the fifteenth of October. Our development plan is sequenced so that "
"a working three-country demonstration is ready well ahead of Round 1, leaving the remaining "
"time for hardening, additional country coverage, and the empirical benchmarking we will report "
"in the memo.")


# ============================================================================
# SLIDE 10 — Closing
# ============================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 0, Inches(0.35), SH, ACCENT)
tb, tf = textbox(s, Inches(0.9), Inches(1.6), Inches(11.5), Inches(0.5))
setpara(tf.paragraphs[0], "IN SUMMARY", 14, RGBColor(0x9F,0xC4,0xE6), bold=True,
        font=FONT_BODY, space_after=0)
tb, tf = textbox(s, Inches(0.9), Inches(2.1), Inches(11.6), Inches(1.4))
setpara(tf.paragraphs[0], "A Deployable Pipeline That Did Not Exist Before",
        40, WHITE, bold=True, font=FONT_HEAD, space_after=0)

points = [
    "Automates ~80% of the RDTII discover → describe workflow for Pillars 6 and 7.",
    "Model-agnostic by design — open-weight, self-hostable, Apache 2.0.",
    "Article-level extraction with all 6 fields, traceable citations, and confidence scores.",
    "Human-first review and a navigable concept graph as the core contribution.",
]
tb, tf = textbox(s, Inches(0.9), Inches(3.7), Inches(11.2), Inches(2.6))
for j, pt in enumerate(points):
    bullet(tf, pt, size=17, color=RGBColor(0xDD,0xE8,0xF2), space_after=14, first=(j == 0),
           marker_color=RGBColor(0x6F,0xB0,0xE0))

rect(s, Inches(0.9), Inches(6.4), Inches(5), Pt(2), ACCENT)
tb, tf = textbox(s, Inches(0.9), Inches(6.6), Inches(11), Inches(0.6))
setpara(tf.paragraphs[0],
        "Thank you.   ·   Repository to be published under Apache 2.0   ·   "
        "escap-digitaltrade-hackathon@un.org",
        13, RGBColor(0xBF,0xD2,0xE5), font=FONT_BODY, space_after=0)

notes(s,
"In closing, let me restate why we believe this project matters. The primary contribution is "
"not a novel model architecture; it is a complete, deployable, end-to-end pipeline that is "
"purpose-built for the RDTII legal extraction task and that does not currently exist in "
"open-source form.\n\n"
"To summarise the four points on screen: first, the system automates approximately eighty "
"percent of the discover-to-describe workflow for Pillars 6 and 7. Second, it is model-agnostic "
"by design, which makes it open-weight capable, fully self-hostable, and released under Apache "
"2.0. Third, it performs genuine article-level extraction, capturing all six mandatory fields "
"with traceable citations and confidence scores. And fourth, it keeps a human firmly in the "
"loop through a non-technical review interface, while contributing a navigable concept graph "
"that turns isolated provisions into connected knowledge. Thank you very much for your "
"attention; I would be glad to take any questions.")


out = r"C:\Users\Drew\Desktop\rdtii-autoextract\docs\RDTII_AutoExtract.pptx"
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
