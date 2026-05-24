"""In-place update of docs/RDTII_AutoExtract.pptx:

  1. Add "Team Arkova" attribution to Slide 1 and Slide 10.
  2. Refine Slide 6 (Stage 5) to reflect the parallel-derivation refinement:
       graph (subtractive) + FCA tree (algebraic) from one shared seed.
  3. Replace all speaker notes with the compressed 4-minute concept-video script.

Preserves layout, shapes, and styling — only the text inside named TextBoxes and
the notes_text_frame is modified.
"""
import sys, io, copy
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")

from pathlib import Path
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
SRC = ROOT / "docs" / "RDTII_AutoExtract.pptx"
prs = Presentation(str(SRC))


def find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None


def set_text_preserving_style(shape, new_text):
    """Overwrite a shape's text, preserving the run formatting of the first run.

    new_text may use ' | ' as a paragraph separator (mirrors how the deck stores
    multi-line bullets in a single TextBox).
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Capture style from the first run we can find
    template_run = None
    for p in tf.paragraphs:
        for r in p.runs:
            template_run = r
            break
        if template_run:
            break
    # Clear all paragraphs except keep the first paragraph element for re-use
    # python-pptx: clear by removing extra <a:p> children
    txBody = tf._txBody
    # Remove all <a:p> children
    for p in list(txBody.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}p")):
        txBody.remove(p)
    lines = new_text.split(" | ")
    for i, line in enumerate(lines):
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        if template_run is not None:
            # Copy rPr (run properties) so font/size/color/bold survive
            from lxml import etree
            src_rPr = template_run._r.find("{http://schemas.openxmlformats.org/drawingml/2006/main}rPr")
            if src_rPr is not None:
                # Remove any auto-created rPr on the new run, then attach a copy
                new_rPr = r._r.find("{http://schemas.openxmlformats.org/drawingml/2006/main}rPr")
                if new_rPr is not None:
                    r._r.remove(new_rPr)
                r._r.insert(0, copy.deepcopy(src_rPr))


def set_notes(slide, text):
    nf = slide.notes_slide.notes_text_frame
    nf.text = text


# ============== CONTENT UPDATES ==============

# --- Slide 1: title — add Team Arkova ---
s1 = prs.slides[0]
# TextBox 7 currently says "Licensed under Apache 2.0". Promote it to attribution line.
set_text_preserving_style(find_shape(s1, "TextBox 7"),
                          "Team Arkova   ·   Licensed under Apache 2.0")

# --- Slide 6: Stage 5 — refine bullets to express parallel derivation ---
s6 = prs.slides[5]
set_text_preserving_style(
    find_shape(s6, "TextBox 8"),
    "—  One shared seed: tagged nodes from Stage 3. Stage 5 derives two artifacts in parallel."
    " | —  Concept graph (subtractive): start from the complete pairwise graph, then drop edges"
    " where w = α·IDF-Jaccard(tags) + (1−α)·cosine(emb) falls below a calibrated θ. Louvain"
    " communities + weighted PageRank score influence."
    " | —  Generality tree (algebraic): Formal Concept Analysis on the node × tag matrix —"
    " independent of graph edges. Intent = tags, extent = sections; more tags ⇒ deeper, more"
    " specific. Together: entry-point category → PageRank-ranked subcategories."
    " | —  Pseudo-deterministic: fixed seeds + θ ⇒ reproducible. Every edge stores its basis"
    " (shared tags, cosine) — no black-box connections."
)

# --- Slide 10: close — add Team Arkova line ---
s10 = prs.slides[9]
set_text_preserving_style(find_shape(s10, "TextBox 7"),
                          "Team Arkova   ·   Repository to be published under Apache 2.0"
                          "   ·   escap-digitaltrade-hackathon@un.org")

# ============== SPEAKER NOTES — 4-MINUTE COMPRESSED SCRIPT ==============
# Total target ~560 words at ~140 wpm. Each block is the spoken narration timed to
# the slide. The hackathon focus on "Automated discovery" and "Mapping and verification"
# gets the dominant share of time.

NOTES = [
    # Slide 1 — 0:00-0:15  (~32 words)
    "[0:00–0:15] We are Team Arkova. We built RDTII AutoExtract — an open-source, "
    "model-agnostic pipeline that automates roughly eighty percent of the RDTII workflow "
    "for Pillars 6 and 7 across Asia-Pacific jurisdictions.",

    # Slide 2 — 0:15-0:40  (~55 words)
    "[0:15–0:40] The problem we tackle is simple: today the RDTII workflow is almost "
    "entirely manual. ESCAP researchers search government portals, retrieve PDFs, read "
    "dense legal text, and extract structured fields — article by article, across many "
    "jurisdictions and languages. It does not scale, and existing literature targets the "
    "EU AI Act and GDPR, not RDTII.",

    # Slide 3 — 0:40-1:00  (~50 words)
    "[0:40–1:00] Our six objectives follow directly. Automate discovery beyond the "
    "ESCAP database. Extract all six mandatory fields at article-level. Map every "
    "provision to Pillar 6 or 7 sub-indicators with verifiable citations. Ship a "
    "self-hostable open tool. Keep every AI component swappable. And provide a "
    "human-first review surface.",

    # Slide 4 — 1:00-1:25  (~55 words)
    "[1:00–1:25] To make swappability real, we build on ports and adapters. The core "
    "domain depends only on interfaces — never on a concrete model. Every AI component "
    "— LLM, OCR, captioner, embeddings, vector store, graph library — is an adapter "
    "behind a port, changed via configuration. No vendor lock-in, no production "
    "proprietary dependency.",

    # Slide 5 — 1:25-2:35  (~155 words)  ←  Automated discovery + start of mapping
    "[1:25–2:35] The pipeline runs in five stages. Stages one and two handle "
    "automated discovery of legal documents. A crawler — Playwright with anti-bot "
    "handling and archive fallback — locates regulations on national portals, gazettes, "
    "and ministry sites, recording full provenance per document. Then OCR converts every "
    "format to clean text at under five percent character error rate, and a "
    "vision-language captioner produces a short descriptor of each section. That caption "
    "becomes the basis for every downstream tag, so we are not hostage to OCR noise on "
    "scanned or non-English documents. "
    "Stage three handles mapping. We tag each section multi-label against indicator "
    "labels, then a RAG-based LLM call extracts the six mandatory fields with strict "
    "JSON schema and a source citation. Stage four is human verification — a "
    "non-technical reviewer console where any mapping is accepted, rejected, or edited "
    "in seconds, with a direct link back to the source span.",

    # Slide 6 — 2:35-3:15  (~95 words)  ←  Concept graph
    "[2:35–3:15] Stage five is our core contribution. From the tagged nodes, we derive "
    "two artifacts in parallel. The concept graph is subtractive — we start from the "
    "complete pairwise graph and drop edges below a calibrated threshold, leaving only "
    "topical borders. Then Louvain communities and weighted PageRank score importance "
    "within the surviving web. Independently, Formal Concept Analysis on the node-by-tag "
    "matrix produces a generality tree: more tags imply fewer matching sections, so "
    "depth equals specificity. Together they give us entry-point categories opening into "
    "ranked subcategories — navigable cross-jurisdiction evidence.",

    # Slide 7 — 3:15-3:35  (~45 words)
    "[3:15–3:35] Every model named here is a starting suggestion, not a commitment. "
    "Development uses Claude or GPT; production resolves to open-weight Llama 3.1 via "
    "Ollama with a single config change. All reused components are Apache-2.0 "
    "compatible, audited for license discipline.",

    # Slide 8 — 3:35-3:50  (~40 words)
    "[3:35–3:50] On viability — open-weight self-hosted cost is roughly five to ten "
    "cents per fifty-page document, API at ten to twenty-five cents. We are "
    "finale-ready for ten countries with no retraining. Everything ships via "
    "docker-compose on commodity hardware.",

    # Slide 9 — 3:50-3:55  (~20 words)  — optional, very brief
    "[3:50–3:55] The path to October runs from the May application through training "
    "workshops, Round 1, the hybrid pitch, and the Bangkok finale.",

    # Slide 10 — 3:55-4:00  (~35 words)
    "[3:55–4:00] To close: not a new model — a complete, deployable, open-source RDTII "
    "pipeline that does not yet exist. Auditable, multilingual, and built for the people "
    "who do the work. Thank you.",
]

for slide, note in zip(prs.slides, NOTES):
    set_notes(slide, note)

# Word count check
total_words = sum(len(n.split()) for n in NOTES)
print(f"Total script words: {total_words} (target ~560 for 4 minutes at 140 wpm)")
for i, n in enumerate(NOTES, 1):
    print(f"  Slide {i:>2}: {len(n.split()):>3} words")

prs.save(str(SRC))
print(f"\nupdated {SRC}")
