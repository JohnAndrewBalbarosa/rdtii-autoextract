"""Dump structure of docs/RDTII_AutoExtract.pptx so we can update in place."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

P = Path(__file__).resolve().parents[2] / "docs" / "RDTII_AutoExtract.pptx"
prs = Presentation(str(P))
print(f"== {P.name} ==")
print(f"Slide size: {prs.slide_width/914400:.2f} x {prs.slide_height/914400:.2f} inches")
print(f"Slides: {len(prs.slides)}")
print(f"Layouts in master: {[l.name for l in prs.slide_layouts]}")
print()
for i, s in enumerate(prs.slides, 1):
    print(f"--- Slide {i} (layout: {s.slide_layout.name}) ---")
    for sh in s.shapes:
        kind = sh.shape_type
        name = sh.name
        if sh.has_text_frame:
            txt = " | ".join(p.text for p in sh.text_frame.paragraphs if p.text)
            if txt.strip():
                print(f"  [{kind}] {name}: {txt[:220]}")
        elif sh.shape_type == 13:  # picture
            print(f"  [PICTURE] {name}")
        else:
            print(f"  [{kind}] {name}")
    if s.has_notes_slide:
        notes = s.notes_slide.notes_text_frame.text
        if notes.strip():
            print(f"  NOTES: {notes[:300]}")
    print()
