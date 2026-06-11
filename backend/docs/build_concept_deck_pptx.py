"""Build CONCEPT_DECK.pptx — concept-video companion deck for Team Arkova.

12 slides, 16:9. Designed to pair with a 4-minute concept video; each slide carries a
short headline + body bullets + speaker-notes timing hint. Mirrors PROPOSAL.md +
TECHNICAL_MEMO.md + GRAPH_PIPELINE.md (parallel derivation in Stage 5).
"""
from pathlib import Path
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = Path(__file__).parent
OUT = HERE / "CONCEPT_DECK.pptx"
if len(sys.argv) > 1:
    OUT = HERE / sys.argv[1]
PIPELINE_PNG = HERE / "pipeline.png"

NAVY = RGBColor(0x0B, 0x3D, 0x91)
BLUE = RGBColor(0x1A, 0x73, 0xE8)
DARK = RGBColor(0x20, 0x20, 0x20)
GREY = RGBColor(0x5F, 0x63, 0x68)
LIGHT = RGBColor(0xF1, 0xF3, 0xF4)
ACCENT = RGBColor(0xE8, 0x71, 0x0A)
GREEN = RGBColor(0x13, 0x73, 0x33)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, size=16, color=DARK):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"•  {item}"
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def header_bar(slide, kicker, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_text(slide, 0.5, 0.07, 12.5, 0.35, kicker, size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, 0.5, 0.7, 12.5, 0.7, title, size=28, bold=True, color=NAVY)


def speaker_notes(slide, timing, body):
    nf = slide.notes_slide.notes_text_frame
    nf.text = f"[{timing}]  {body}"


def new_slide():
    return prs.slides.add_slide(BLANK)


# ============== SLIDE 1 — Title ==============
s = new_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
add_text(s, 0.7, 2.4, 12, 1.2, "Zetarix", size=54, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, 0.7, 3.6, 12, 0.7, "Open-Source AI Pipeline for Digital Trade Regulatory Analysis", size=22, color=RGBColor(0xE8, 0xF0, 0xFE), italic=True)
add_text(s, 0.7, 4.5, 12, 0.5, "Pillar 6 · Cross-Border Data Flows    |    Pillar 7 · Domestic Data Protection", size=16, color=RGBColor(0xE8, 0xF0, 0xFE))
add_text(s, 0.7, 6.4, 12, 0.5, "Team Arkova", size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, 0.7, 6.85, 12, 0.4, "UN ESCAP & KMITL Global Hackathon · 2026 · Apache 2.0", size=12, color=RGBColor(0xE8, 0xF0, 0xFE))
speaker_notes(s, "0:00–0:15", "Open with the project name and our scope: Pillars 6 and 7, RDTII, Team Arkova.")

# ============== SLIDE 2 — Problem ==============
s = new_slide()
header_bar(s, "THE PROBLEM", "Manual RDTII workflow does not scale")
add_bullets(s, 0.7, 1.8, 12, 4.5, [
    "ESCAP researchers manually search portals, retrieve PDFs, and read dense legal text.",
    "Per-article extraction across multiple jurisdictions — bottlenecked by researcher time and language coverage.",
    "Asia-Pacific scope spans English, Vietnamese, Thai, Bahasa Malaysia, and more — scanned PDFs and messy portals are routine.",
    "Existing literature targets GDPR / the EU AI Act — no open-source pipeline targets the RDTII indicators end-to-end.",
])
speaker_notes(s, "0:15–0:45", "Frame the burden: manual workflow, multilingual, no existing open-source solution targets RDTII directly.")

# ============== SLIDE 3 — Thesis ==============
s = new_slide()
header_bar(s, "OUR THESIS", "Automate 80% — keep 20% transparent human review")
add_bullets(s, 0.7, 1.8, 12, 4.5, [
    "Automate discovery → describe (Stages 1–4); leave verification to a non-technical reviewer console.",
    "Article-level granularity — every output is a section with the 6 mandatory RDTII fields.",
    "Model-agnostic by design: every AI component sits behind a port; concrete models are swappable adapters.",
    "Self-hostable via docker-compose — no proprietary production dependency.",
])
add_text(s, 0.7, 6.6, 12, 0.4, "title · last_update · url · scope · provisions · impact", size=14, italic=True, color=GREY)
speaker_notes(s, "0:45–1:10", "Our thesis: 80/20 split, article-level, model-agnostic, self-hostable. The 6 fields are the contract.")

# ============== SLIDE 4 — Pipeline overview (with PNG) ==============
s = new_slide()
header_bar(s, "SYSTEM ARCHITECTURE", "Five-stage pipeline behind ports & adapters")
if PIPELINE_PNG.exists():
    s.shapes.add_picture(str(PIPELINE_PNG), Inches(0.4), Inches(1.7), width=Inches(12.5))
add_text(s, 0.7, 6.55, 12, 0.7,
         "Discover  →  OCR + Caption  →  Tag + Extract  →  Human Review  →  Concept Graph",
         size=14, italic=True, color=GREY, align=PP_ALIGN.CENTER)
speaker_notes(s, "1:10–1:35", "Walk the audience through the five stages at a high level. Dashed boxes = the ports each stage binds to.")

# ============== SLIDE 5 — Stage 1+2 (Automated discovery) ==============
s = new_slide()
header_bar(s, "AUTOMATED DISCOVERY  ·  STAGES 1–2", "Crawl, OCR, and caption — the basis for every tag")
add_text(s, 0.7, 1.7, 6, 0.5, "Stage 1 — Discover", size=18, bold=True, color=BLUE)
add_bullets(s, 0.7, 2.2, 6, 3.5, [
    "Playwright + BeautifulSoup crawler",
    "Compliant access handling; archive fallback",
    "Per-document provenance (URL, timestamp, country, language, source)",
    "Language detection routes non-English to translation",
], size=14)
add_text(s, 7.0, 1.7, 6, 0.5, "Stage 2 — OCR + Caption", size=18, bold=True, color=BLUE)
add_bullets(s, 7.0, 2.2, 6, 3.5, [
    "Tesseract 5 + OpenCV; PaddleOCR fallback for SE-Asian scripts",
    "Target < 5% character error rate",
    "Vision-language captioning (Qwen2-VL / ColQwen2) on scanned or non-English docs",
    "The caption is the BASIS for every downstream tag — garbage caption ⇒ garbage tags",
], size=14)
speaker_notes(s, "1:35–2:15", "Explain that captioning happens up front and is what tagging scores against — not raw OCR bytes.")

# ============== SLIDE 6 — Stage 3 (Tag + Extract) ==============
s = new_slide()
header_bar(s, "MAPPING  ·  STAGE 3", "Tag, extract, and map to RDTII indicators")
add_bullets(s, 0.7, 1.8, 12, 4.5, [
    "Multi-label tagging (BGE-M3 embeddings; sigmoid scoring of caption against indicator labels — a section can carry P6 and P7).",
    "RAG over pgvector / Chroma → LLM with strict JSON schema enforcing all 6 fields + citation.",
    "Per-record confidence score; low-confidence rows surface first in review.",
    "LLM provider is a port — dev uses Claude / GPT, production swaps to Llama 3.1 via Ollama with no code change.",
], size=15)
speaker_notes(s, "2:15–2:45", "Show that extraction outputs are structured (JSON), cited, scored — and the LLM is interchangeable.")

# ============== SLIDE 7 — Stage 4 (Human Review console) ==============
s = new_slide()
header_bar(s, "VERIFICATION  ·  STAGE 4", "The reviewer console — the final 20% as a first-class surface")
add_bullets(s, 0.7, 1.8, 12, 4.5, [
    "Each row = one extracted article with all 6 fields inline; expand for scope / provisions / impact.",
    "Per-row Verify / Reject / Reset — every action immutable, audit-tracked.",
    "Pillar tag, indicator code, confidence meter, source language badge, and link to the exact source span.",
    "Built for non-technical ESCAP researchers — no JSON, no code. Filter by status, pillar, and free-text search.",
], size=15)
speaker_notes(s, "2:45–3:10", "Reviewer console is the user-facing surface — not an afterthought. Verification is fast.")

# ============== SLIDE 8 — Stage 5 (Concept Graph) ==============
s = new_slide()
header_bar(s, "CONCEPT GRAPH  ·  STAGE 5 (CORE CONTRIBUTION)", "One seed, two parallel artifacts")
add_text(s, 0.7, 1.7, 12, 0.5,
         "Tagged nodes are the shared seed — the graph (subtractive) and the tree (algebraic) are derived in parallel.",
         size=14, italic=True, color=GREY)
# Two columns
add_text(s, 0.7, 2.5, 6, 0.5, "Concept Graph (subtractive)", size=18, bold=True, color=BLUE)
add_bullets(s, 0.7, 3.0, 6, 3.5, [
    "Start from the complete pairwise graph",
    "Edge weight = α·IDF-Jaccard(tags) + (1−α)·cosine(emb)",
    "Drop edges below calibrated θ (fixed seeds ⇒ reproducible)",
    "Louvain communities + weighted PageRank",
], size=14)
add_text(s, 7.0, 2.5, 6, 0.5, "FCA Tree (algebraic)", size=18, bold=True, color=GREEN)
add_bullets(s, 7.0, 3.0, 6, 3.5, [
    "Built directly from the node × tag matrix — independent of graph edges",
    "Intent = tags  ·  Extent = sections",
    "More tags ⇒ fewer matching sections ⇒ deeper / more specific",
    "Entry-point categories at the root; ranked subcategories below",
], size=14)
speaker_notes(s, "3:10–3:40", "This is the headline contribution: the tree is NOT derived from the graph. Two artifacts, one seed.")

# ============== SLIDE 9 — Ports & Adapters table ==============
s = new_slide()
header_bar(s, "MODEL-AGNOSTIC BY DESIGN", "Every component is a swappable adapter")
rows = [
    ("Port", "Dev adapter", "Open-weight target", "License"),
    ("DocumentSource", "Playwright + BS4", "same", "BSD/MIT"),
    ("OCREngine", "Tesseract 5 / PaddleOCR", "same", "Apache-2.0"),
    ("Captioner", "Qwen2-VL / ColQwen2", "same", "Apache-2.0"),
    ("Tagger", "BGE-M3 + reranker", "same", "MIT/Apache"),
    ("VectorStore", "pgvector / Chroma", "same", "PostgreSQL/Apache"),
    ("LLMProvider", "Claude / GPT", "Llama 3.1 8B/70B (Ollama)", "Apache-compat"),
    ("GraphBuilder / Ranker", "NetworkX + FCA `concepts`", "same", "BSD/MIT"),
]
left, top = 0.5, 1.7
col_w = [2.8, 3.2, 3.6, 2.6]
row_h = 0.42
table_shape = s.shapes.add_table(len(rows), len(rows[0]),
                                  Inches(left), Inches(top),
                                  Inches(sum(col_w)), Inches(row_h * len(rows)))
table = table_shape.table
for j, w in enumerate(col_w):
    table.columns[j].width = Inches(w)
for i, row in enumerate(rows):
    for j, text in enumerate(row):
        cell = table.cell(i, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = text
        r.font.name = "Calibri"
        r.font.size = Pt(12)
        r.font.bold = (i == 0)
        if i == 0:
            r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT if i % 2 == 0 else RGBColor(0xFF, 0xFF, 0xFF)
add_text(s, 0.5, 6.55, 12, 0.5,
         "License discipline: Apache-2.0-compatible only — ColQwen2 over PaliGemma, NetworkX over GPL leidenalg, no CC-BY-NC models.",
         size=11, italic=True, color=GREY)
speaker_notes(s, "3:40–4:00", "Every model named is a starting suggestion, not a commitment. Production resolves to open-weight.")

# ============== SLIDE 10 — Cost + generalisation ==============
s = new_slide()
header_bar(s, "COST & GENERALISATION", "Cheap to run · ready for unseen jurisdictions")
add_text(s, 0.7, 1.7, 6, 0.5, "Cost / 50-page document", size=18, bold=True, color=BLUE)
add_bullets(s, 0.7, 2.2, 6, 2.5, [
    "Open-weight (Llama 3.1 8B, self-host GPU):  ~USD 0.05–0.10  (API: $0)",
    "API (Claude / GPT):  ~USD 0.10–0.25",
    "Final benchmarks reported after empirical testing.",
], size=14)
add_text(s, 7.0, 1.7, 6, 0.5, "Generalisation", size=18, bold=True, color=BLUE)
add_bullets(s, 7.0, 2.2, 6, 3.5, [
    "No per-country retraining",
    "Non-English via captioning + translation",
    "Scanned PDFs via OCR (<5% CER)",
    "Messy portals via headless browser + archive fallback",
    "Target ≥ 3 of 10 assigned countries at finale",
], size=14)
speaker_notes(s, "—", "Off-script slide for Q&A; do not need to land on it in the 4-minute version.")

# ============== SLIDE 11 — Originality ==============
s = new_slide()
header_bar(s, "ORIGINALITY", "Not a new model — a deployable RDTII pipeline that does not yet exist")
add_bullets(s, 0.7, 1.8, 12, 4.5, [
    "Closest analogues (e.g. EU AI-Act extraction with LLM + knowledge graphs, 2025) target Western instruments — not RDTII.",
    "Our contribution: end-to-end open-source pipeline + reviewer console + concept-graph evidence layer specifically for RDTII Pillars 6/7 across APAC.",
    "Fine-tuning planned: few-shot tune a small encoder/classifier for sub-indicator mapping; θ calibrated F1-optimal on RDTII labelled data — not guessed.",
    "All fine-tuned weights, training data, and methodology disclosed in the Technical Memo and published to the GitHub repo.",
], size=14)
speaker_notes(s, "—", "Originality framing for judges: novel composition + novel target, not a novel model.")

# ============== SLIDE 12 — Close ==============
s = new_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
add_text(s, 0.7, 2.6, 12, 1.0, "An auditable, open, self-hostable RDTII engine.", size=36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, 0.7, 3.8, 12, 0.7, "docker-compose up  →  api · web · postgres+pgvector", size=20, color=RGBColor(0xE8, 0xF0, 0xFE), italic=True, font="Consolas")
add_text(s, 0.7, 5.0, 12, 0.5, "Built by Team Arkova", size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, 0.7, 5.5, 12, 0.4, "UN ESCAP & KMITL Global Hackathon on AI for Digital Trade Regulatory Analysis · 2026 · Apache 2.0",
         size=12, color=RGBColor(0xE8, 0xF0, 0xFE))
speaker_notes(s, "3:50–4:00", "Land the closing line: deployable, auditable, open. Cue to questions or video CTA.")

prs.save(OUT)
print(f"wrote {OUT}")
